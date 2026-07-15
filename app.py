from __future__ import annotations

import os
import re
import secrets
import time
import json
import html
from io import BytesIO
from pathlib import Path
from typing import Any
from urllib import request as urllib_request
from urllib.error import URLError

from docx import Document
from flask import Flask, Response, abort, jsonify, render_template, request, send_file, url_for
from markupsafe import Markup
from pptx import Presentation
from pypdf import PdfReader
from werkzeug.datastructures import FileStorage
from werkzeug.utils import secure_filename

from raw_materials.chunker import chunk_text
from raw_materials.jsonl_builder import (
    SOURCE_TYPE_TO_MENTOR_MODE,
    build_records,
    default_mentor_mode,
    records_to_jsonl,
    records_to_preview_markdown,
    records_to_pretty_json,
    slugify,
)
from raw_materials.prompt_builder import (
    MODE_DEFINITIONS,
    PromptArtifact,
    build_combined_prompt_text,
    build_mode_prompt_artifacts,
)
from raw_materials.reader import SUPPORTED_EXTENSIONS, extract_text_from_upload, is_supported_filename

try:
    import requests
except ImportError:  # pragma: no cover - used only when optional dependency is absent
    class _CompatRequestException(Exception):
        pass

    class _CompatResponse:
        def __init__(self, body: bytes, status: int):
            self._body = body
            self.status_code = status

        def raise_for_status(self) -> None:
            if self.status_code >= 400:
                raise _CompatRequestException(f"HTTP {self.status_code}")

        def json(self) -> dict[str, Any]:
            return json.loads(self._body.decode("utf-8"))

    class _RequestsCompat:
        RequestException = _CompatRequestException

        @staticmethod
        def post(url: str, json: dict[str, str], headers: dict[str, str], timeout: int) -> _CompatResponse:
            data = __import__("json").dumps(json).encode("utf-8")
            req = urllib_request.Request(url, data=data, headers=headers, method="POST")
            try:
                with urllib_request.urlopen(req, timeout=timeout) as response:
                    return _CompatResponse(response.read(), response.status)
            except URLError as exc:
                raise _CompatRequestException(str(exc)) from exc

    requests = _RequestsCompat()

try:
    from dotenv import load_dotenv
except ImportError:  # pragma: no cover - optional local convenience dependency
    def load_dotenv() -> None:
        return None


load_dotenv()

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 20 * 1024 * 1024
app.config["OUTPUT_DIR"] = Path(__file__).parent / "outputs"
app.config["MENTOR_LIBRARY_DIR"] = Path(__file__).parent / "mentor_files"

MAX_PROMPT_LENGTH = 4_000
MAX_EXTRACTED_TEXT = 100_000
MAX_OLLAMA_REFERENCE_CHARS = 12_000
MAX_PROMPT_PREVIEW_SEGMENTS = 3

UPLOAD_TYPES = {
    "document": {
        "label": "Document",
        "extensions": {".pdf", ".docx", ".txt"},
        "formats": "PDF · DOCX · TXT",
        "accept": ".pdf,.docx,.txt",
        "description": "Reports, proposals, essays, plans, and other written work.",
        "focus": "e.g. clarity and structure",
    },
    "transcript": {
        "label": "Transcript",
        "extensions": {".txt", ".srt", ".vtt", ".docx"},
        "formats": "TXT · SRT · VTT · DOCX",
        "accept": ".txt,.srt,.vtt,.docx",
        "description": "Meetings, interviews, conversations, captions, and recorded sessions.",
        "focus": "e.g. key themes and actions",
    },
    "powerpoint": {
        "label": "PowerPoint",
        "extensions": {".pptx"},
        "formats": "PPTX",
        "accept": ".pptx",
        "description": "Presentation slides, speaker notes, pitch decks, and briefings.",
        "focus": "e.g. story flow and messaging",
    },
}

MENTORS = {
    "dr-nanshu-lu": {
        "name": "Dr. Nanshu Lu",
        "initials": "NL",
        "status": "Available mentor",
        "description": (
            "Her specialty, thinking process, and feedback style are configured "
            "in the local model's mentor profile."
        ),
    }
}
DEFAULT_MENTOR_ID = "dr-nanshu-lu"

DOWNLOAD_FILENAMES = {
    "jsonl": "records.jsonl",
    "json": "records_pretty.json",
    "md": "preview.md",
}

REFERENCE_UPLOAD_GROUPS = {
    "meeting_research_pi": {
        "field": "research_files",
        "label": "Research Ideas / Meeting Minutes",
        "description": "Meeting notes, research ideas, experiment plans, and lab discussion records.",
    },
    "slides_talk_pi": {
        "field": "slide_files",
        "label": "Talks / Presentations / Slides",
        "description": "Talk drafts, presentation slides, figure sets, and slide feedback.",
    },
    "paper_proposal_pi": {
        "field": "paper_files",
        "label": "Papers / Proposals",
        "description": "Manuscripts, proposals, paper drafts, reviewer comments, and cover letters.",
    },
}


def call_model(
    prompt: str,
    demo_feedback: str | None = None,
    mentor_id: str | None = None,
) -> str:
    """Send a prompt to the configured model, or return local demo feedback."""
    model_url = os.getenv("MODEL_API_URL", "").strip()
    if not model_url:
        time.sleep(0.4)
        return demo_feedback or (
            "This is a demo response. Your Python website is working. Add the "
            "local model address to .env as MODEL_API_URL when it is ready."
        )

    headers = {"Content-Type": "application/json"}
    api_key = os.getenv("MODEL_API_KEY", "").strip()
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    payload: dict[str, str] = {"prompt": prompt}
    if mentor_id in MENTORS:
        payload["mentor_id"] = mentor_id or ""
        payload["mentor_name"] = MENTORS[mentor_id]["name"]

    response = requests.post(model_url, json=payload, headers=headers, timeout=120)
    response.raise_for_status()
    data = response.json()
    for key in ("output", "response", "text"):
        value = data.get(key)
        if isinstance(value, str) and value.strip():
            return value.strip()
    raise ValueError("Model response did not include output, response, or text.")


def call_ollama_generate(prompt: str) -> str:
    """Generate text with a local Ollama model."""
    base_url = os.getenv("OLLAMA_URL", "http://127.0.0.1:11434").strip().rstrip("/")
    model = os.getenv("OLLAMA_MODEL", "llama3.1:8b").strip() or "llama3.1:8b"
    response = requests.post(
        f"{base_url}/api/generate",
        json={
            "model": model,
            "prompt": prompt,
            "stream": False,
            "options": {
                "temperature": 0.2,
                "top_p": 0.9,
            },
        },
        headers={"Content-Type": "application/json"},
        timeout=180,
    )
    response.raise_for_status()
    data = response.json()
    value = data.get("response")
    if not isinstance(value, str) or not value.strip():
        raise ValueError("Ollama response did not include a non-empty response field.")
    return value.strip()


def prompt_llm_provider() -> str:
    return os.getenv("PROMPT_LLM_PROVIDER", "").strip().lower()


def extract_generated_prompt(content: str) -> str:
    marker = "Generated PI-style prompt:"
    if marker not in content:
        return content.strip()
    after_marker = content.split(marker, 1)[1]
    return after_marker.split("PI-style response rules:", 1)[0].strip()


def build_ollama_style_distillation_request(
    mode: str,
    source_files: list[str],
    chunks: list[str],
) -> str:
    """Ask Ollama to extract reusable PI review moves before writing a prompt."""
    definition = MODE_DEFINITIONS[mode]
    reference_text = "\n\n".join(chunks)[:MAX_OLLAMA_REFERENCE_CHARS]
    source_line = "; ".join(source_files) if source_files else "uploaded reference materials"
    if mode == "slides_talk_pi":
        extraction_guidance = (
            "Extract reusable slide/talk review moves in this exact spirit:\n"
            "- audience comprehension: how the PI makes the talk easy to follow through labels, callouts, and visual cues\n"
            "- titles and significance framing: how titles, opening context, and background slides communicate broad importance without narrowing impact\n"
            "- citation discipline: how every non-original figure, chronology, comparison, or claim is cited in the right place\n"
            "- visual consistency: how similar plots, colors, labels, legends, and formats stay consistent across slides\n"
            "- takeaway messages: how each slide or section tells the audience what to remember\n"
            "- remove weak or amateur-looking visuals: how the PI deletes poor curves, decorative figures, clutter, or unprofessional illustrations\n"
            "- slide-level fixes: how the PI gives concrete layout, color, label, citation, and figure-placement edits\n\n"
        )
    elif mode == "paper_proposal_pi":
        extraction_guidance = (
            "Extract reusable manuscript/proposal review moves in this exact spirit:\n"
            "- practical value: how the title, abstract, and opening explain why the work matters\n"
            "- coherent argument: how literature, figures, discussion, and conclusion reinforce one central claim\n"
            "- boundaries: how applicable range, setup, controls, and usage guidance are clarified\n"
            "- context: how literature, materials comparisons, standards, citations, and applications broaden the frame\n"
            "- figure proof: how figures prove the main claim through motivation, mechanism, comparisons, transitions, and decisive evidence\n"
            "- precision: how terminology, labels, captions, legends, axes, colors, and layout are made publication-ready\n"
            "- concrete revisions: how feedback turns into edits to title, figures, captions, discussion, references, or structure\n\n"
        )
    else:
        extraction_guidance = (
            "Extract reusable review moves in this exact spirit:\n"
            "- reframe: how the PI reframes a technical task as a broader research opportunity\n"
            "- ground: how the PI grounds the idea in standards, literature, user needs, industry examples, or practical context\n"
            "- decompose: how the PI breaks the problem into variables, design parameters, constraints, and measurable outcomes\n"
            "- compare mechanisms: how the PI asks for competing mechanisms, alternative explanations, and necessary comparisons\n"
            "- prioritize: how the PI separates immediate roadmap priorities from later variables or side projects\n"
            "- action items: how the PI turns discussion into papers to read, data to collect, collaborators to contact, or experiments to run\n\n"
        )
    return (
        "You are distilling a professor's review style from raw reference materials.\n"
        "Do not write the final reusable prompt yet. Extract reusable review moves only.\n"
        "The uploaded project is an example for style distillation, not the future review target.\n"
        "Do not copy project-specific nouns, material systems, sample names, mechanisms, or project goals.\n\n"
        f"Mode: {definition['label']}\n"
        f"Source files: {source_line}\n\n"
        f"{extraction_guidance}"
        "Return 5-7 concise bullets. Each bullet must describe a reusable review habit, not a project detail.\n\n"
        f"Uploaded reference material excerpts:\n{reference_text}"
    )


def build_ollama_prompt_request(
    mode: str,
    source_files: list[str],
    chunks: list[str],
    deterministic_prompt: str,
    distilled_pattern: str | None = None,
) -> str:
    definition = MODE_DEFINITIONS[mode]
    source_line = "; ".join(source_files) if source_files else "uploaded reference materials"
    pattern = distilled_pattern.strip() if distilled_pattern and distilled_pattern.strip() else (
        "Use the mode priorities and deterministic draft as the style pattern."
    )
    rhetorical_skeleton = ""
    mode_specific_guidance = ""
    if mode == "meeting_research_pi":
        mode_specific_guidance = (
            "The final paragraph must preserve the concrete advisor moves when they are present in the pattern: "
            "reframe the research opportunity, ground it in practical context such as standards/literature/user needs/industry, "
            "decompose the problem into variables or design parameters, compare mechanisms or competing explanations, "
            "prioritize roadmap and scope, and convert the discussion into concrete action items. "
            "Do not collapse everything into a generic hypothesis-controls-evidence checklist; hypothesis, controls, "
            "and evidence are useful only when integrated with framing, mechanism comparison, prioritization, and next actions.\n\n"
            "A strong final paragraph should use concrete verbs like reframe, ground, decompose, compare, prioritize, "
            "and translate into action items when the distilled pattern supports them.\n\n"
            "For research ideas or meeting minutes, the final prompt should explicitly cover this six-move sequence: "
            "reframe -> ground -> decompose -> compare -> prioritize -> action items. "
            "Do not let falsifiability, controls, or evidence dominate the paragraph; include them only as part of the "
            "broader advisor workflow. Use action verbs. Make the final clause emphasize concrete next steps such as "
            "papers to read, data to collect, collaborators to contact, comparisons to run, or experiments to prioritize.\n\n"
        )
        rhetorical_skeleton = (
            "Use this rhetorical skeleton for the final paragraph, adapting wording but preserving the logic: "
            "I first ask whether the idea has been reframed from a technical task into a broader research opportunity "
            "with clear practical relevance. The discussion should be grounded in real context, such as standards, "
            "literature, industry examples, user needs, or measurable pain points. I then look for a clean decomposition "
            "of the problem into variables, design parameters, constraints, and measurable outcomes, followed by a "
            "comparison of competing mechanisms or alternative explanations. The feedback should distinguish immediate "
            "roadmap priorities from later variables or side projects before judging detailed experiments, and translate the discussion into concrete action "
            "items: papers to read, data to collect, collaborators to contact, comparisons to run, or experiments to "
            "prioritize. Ultimately, the next experiment should be capable of changing the project direction. "
            "Avoid repeating user needs or any other criterion twice.\n\n"
        )
    elif mode == "slides_talk_pi":
        rhetorical_skeleton = (
            "Use this Talks/Presentations/Slides rhetorical skeleton for the final paragraph, adapting wording but preserving "
            "the logic; avoid a generic design-polish checklist and avoid a chain of question-form sentences. Begin with whether "
            "the audience can follow the story and immediately understand the scientific logic from the title and opening framing, as well as the slide sequence. "
            "Each slide should help the audience understand the need, logic, and takeaway of the work rather than simply display "
            "information. Emphasize specific slide titles, broad significance framing that does not artificially narrow the "
            "significance of the work, and background or chronology slides that establish a clear need for the research. "
            "The prompt should require the reviewer to cite every non-original figure, claim, chronology, or comparison in the "
            "right visual location, and to add labels, annotations, parentheses, callouts, and visual cues wherever they help "
            "the audience follow. Similar concepts should use consistent labels, annotations, colors, legends, and plot formats "
            "across the talk. The prompt should explicitly delete weak, confusing, amateur-looking, or low-information slides. "
            "Weak curves, decorative figures, amateur-looking illustrations, or low-information slides should be removed or replaced with concise summaries. "
            "The review should ultimately translate into concrete slide-level "
            "edits to titles, citations, labels, figure choices, layout fixes, visual consistency, takeaway messages, and "
            "audience guidance.\n\n"
        )
    elif mode == "paper_proposal_pi":
        rhetorical_skeleton = (
            "Use this Papers/Proposals rhetorical skeleton for the final paragraph, adapting wording but preserving the logic; "
            "avoid a checklist-like sequence of repeated 'I expect', 'I require', or 'should' sentences. Begin with whether "
            "the title, abstract, and opening narrative clearly communicate the practical value of the work, why the reader "
            "should care, and what central claim the manuscript is trying to establish. Then evaluate whether the manuscript "
            "builds a coherent argument in which literature positioning, figures, discussion, and conclusion reinforce the same "
            "claim rather than functioning as separate sections. Pay close attention to the applicable range, boundary conditions, "
            "experimental setup, and evidence behind each claim while checking whether the context is properly supported by "
            "relevant literature, materials comparisons, standards, citations, or application scenarios. Examine whether the "
            "figures actually prove the main claim through motivation, mechanism, comparison, transitions, and decisive evidence "
            "rather than decorative or low-information panels. Also evaluate precise terminology and consistent captions, labels, "
            "legends, colors, axes, panel alignment, and layout so that claim-evidence alignment is maintained throughout. "
            "The review should ultimately translate these issues into concrete revisions to the title, figures, captions, "
            "discussion, references, or manuscript structure, including the abstract when needed.\n\n"
        )
    return (
        "You are helping build a reusable PI-style review prompt from raw reference materials.\n"
        "Do not write a prompt for the uploaded project. The uploaded project is only a reference example "
        "for learning the professor's thinking pattern.\n"
        "Task: write ONE polished, professional, copy-ready, general-purpose prompt paragraph for a future project reviewer.\n"
        "Do not summarize the files for the user. Do not mention that you are an AI. Do not use markdown.\n"
        "The paragraph must be fluent, mode-specific, directly useful, and reusable across projects.\n"
        "Do NOT mention project-specific nouns, material systems, sample names, device names, project names, "
        "or mechanisms from the uploaded reference. Generalize project-specific content into broad review habits, "
        "not topic details.\n"
        "Write a single refined paragraph, not a checklist. Do not repeat the same checklist in different words. "
        "Avoid internal or mechanical phrases such as 'for a project in this mode', 'Specifically, test whether', "
        "'this uploaded-material signal', or 'mode priorities'. Before returning, silently revise your draft once "
        "for elegance, concision, non-redundancy, and copy-ready wording.\n"
        "Bad output: a prompt about the uploaded project itself. Good output: a prompt that can review any future project "
        "in the same mode using the learned PI style.\n\n"
        f"Mode: {definition['label']}\n"
        f"Source files: {source_line}\n"
        f"Mode priorities, to include only if they fit the pattern naturally: {definition['priority_sentence']}\n"
        f"Expected response shape, as optional guidance only: {definition['deliverable_sentence']}\n\n"
        "Use the distilled PI review pattern below as the main input. It already abstracts away project details.\n\n"
        f"Distilled PI review pattern:\n{pattern}\n\n"
        f"{mode_specific_guidance}"
        f"{rhetorical_skeleton}"
        f"General deterministic draft to improve:\n{deterministic_prompt}\n\n"
        "Return only the final general-purpose single refined paragraph, 90-140 words."
    )


def generate_llm_prompt_for_mode(
    mode: str,
    source_files: list[str],
    chunks: list[str],
    deterministic_prompt: str,
) -> str | None:
    if prompt_llm_provider() != "ollama" or not chunks:
        return None
    try:
        distilled_pattern = call_ollama_generate(
            build_ollama_style_distillation_request(mode, source_files, chunks)
        )
        prompt = build_ollama_prompt_request(
            mode,
            source_files,
            chunks,
            deterministic_prompt,
            distilled_pattern=distilled_pattern,
        )
        return polish_llm_prompt_output(call_ollama_generate(prompt))
    except (requests.RequestException, ValueError) as exc:
        app.logger.warning("Ollama prompt generation failed for %s: %s", mode, exc)
        return None


def polish_llm_prompt_output(text: str) -> str:
    """Remove common mechanical LLM artifacts from generated prompt paragraphs."""
    polished = text.strip().strip('"').strip("'").strip()
    polished = polished.replace("\n", " ")
    mechanical_phrases = [
        " for a project in this mode",
        " in this mode",
        " this uploaded-material signal",
        " mode priorities",
    ]
    for phrase in mechanical_phrases:
        polished = polished.replace(phrase, "")
        polished = polished.replace(phrase.title(), "")
    polished = re.sub(
        r"\s*Specifically,\s+test whether\b.*?(?:\.\s*|$)",
        " ",
        polished,
        flags=re.IGNORECASE,
    )
    polished = re.sub(r"\s+", " ", polished).strip()
    polished = re.sub(r"\s+([,.;:])", r"\1", polished)
    if polished and polished[-1] not in ".!?":
        polished += "."
    return polished


def extract_plain_text(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="replace")


def extract_docx(file_bytes: bytes) -> str:
    document = Document(BytesIO(file_bytes))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    table_cells: list[str] = []
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    table_cells.append(text)
    return "\n".join(paragraphs + table_cells)


def extract_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(BytesIO(file_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(page.strip() for page in pages if page.strip())


def extract_powerpoint(file_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(file_bytes))
    slide_text: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        pieces: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                pieces.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            pieces.append(text)
        if pieces:
            slide_text.append(f"Slide {index}: " + "\n".join(pieces))
    return "\n\n".join(slide_text)


def extract_feedback_text(file_bytes: bytes, extension: str) -> str:
    extension = extension.lower()
    if extension in {".txt", ".srt", ".vtt"}:
        return extract_plain_text(file_bytes)
    if extension == ".docx":
        return extract_docx(file_bytes)
    if extension == ".pdf":
        return extract_pdf(file_bytes)
    if extension == ".pptx":
        return extract_powerpoint(file_bytes)
    raise ValueError("Unsupported file type.")


def build_feedback_prompt(
    kind: str,
    filename: str,
    content: str,
    focus: str,
    mentor_id: str,
) -> str:
    label = UPLOAD_TYPES[kind]["label"]
    mentor_name = MENTORS[mentor_id]["name"]
    focus_line = f"\nFeedback focus requested by user: {focus.strip()}" if focus.strip() else ""
    clipped_content = content[:MAX_EXTRACTED_TEXT]
    return (
        f"You are {mentor_name}. Use your configured specialty, thinking process, "
        f"and feedback style to review this {label.lower()}.\n"
        f"File name: {filename}{focus_line}\n\n"
        "Give concise, specific, actionable feedback with strengths, concerns, "
        "and next steps.\n\n"
        f"Content:\n{clipped_content}"
    )


def get_output_dir() -> Path:
    return Path(app.config["OUTPUT_DIR"])


def get_mentor_library_dir() -> Path:
    return Path(app.config["MENTOR_LIBRARY_DIR"])


def mentor_slug(name: str) -> str:
    words = re.findall(r"[A-Za-z0-9]+", name.lower())
    return "-".join(words[:8]) if words else "mentor"


def mentor_dir(slug: str) -> Path:
    safe_slug = mentor_slug(slug)
    return get_mentor_library_dir() / safe_slug


def mode_dir_for_mentor(slug: str, mode: str) -> Path:
    return mentor_dir(slug) / mode


def ensure_prompt_mentor(name: str) -> dict[str, str]:
    display_name = name.strip() or "PI Style Library"
    slug = mentor_slug(display_name)
    directory = mentor_dir(slug)
    directory.mkdir(parents=True, exist_ok=True)
    for mode in MODE_DEFINITIONS:
        (directory / mode / "raw").mkdir(parents=True, exist_ok=True)
    metadata = {"slug": slug, "name": display_name}
    (directory / "mentor.json").write_text(json.dumps(metadata, indent=2), encoding="utf-8")
    return metadata


def read_prompt_mentor(slug: str) -> dict[str, str] | None:
    safe_slug = mentor_slug(slug)
    directory = mentor_dir(safe_slug)
    metadata_path = directory / "mentor.json"
    if not metadata_path.exists():
        return None
    try:
        data = json.loads(metadata_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {"slug": safe_slug, "name": safe_slug.replace("-", " ").title()}
    return {
        "slug": str(data.get("slug") or safe_slug),
        "name": str(data.get("name") or safe_slug.replace("-", " ").title()),
    }


def list_prompt_mentors() -> list[dict[str, str]]:
    root = get_mentor_library_dir()
    if not root.exists():
        return []
    mentors: list[dict[str, str]] = []
    for child in sorted(root.iterdir(), key=lambda path: path.name.lower()):
        if not child.is_dir():
            continue
        mentor = read_prompt_mentor(child.name)
        if mentor:
            mentors.append(mentor)
    return mentors


def resolve_prompt_mentor(selected_slug: str = "", new_name: str = "") -> dict[str, str]:
    if new_name.strip():
        return ensure_prompt_mentor(new_name)
    if selected_slug.strip():
        existing = read_prompt_mentor(selected_slug)
        if existing:
            return existing
        raise ValueError("Please select an existing mentor or create a new one.")
    return ensure_prompt_mentor("PI Style Library")


def stored_files_for_mentor(slug: str) -> dict[str, list[str]]:
    if not slug:
        return {mode: [] for mode in MODE_DEFINITIONS}
    files_by_mode: dict[str, list[str]] = {}
    for mode in MODE_DEFINITIONS:
        raw_dir = mode_dir_for_mentor(slug, mode) / "raw"
        files_by_mode[mode] = (
            sorted(path.name for path in raw_dir.iterdir() if path.is_file())
            if raw_dir.exists()
            else []
        )
    return files_by_mode


def save_uploaded_reference_files(slug: str) -> None:
    for mode, config in REFERENCE_UPLOAD_GROUPS.items():
        raw_dir = mode_dir_for_mentor(slug, mode) / "raw"
        raw_dir.mkdir(parents=True, exist_ok=True)
        files = [file for file in request.files.getlist(config["field"]) if file and file.filename]
        for uploaded_file in files:
            filename = safe_uploaded_filename(uploaded_file)
            if not is_supported_filename(filename):
                raise ValueError(
                    f"Unsupported file type for {config['label']}: {filename}. "
                    f"Use one of: {', '.join(sorted(SUPPORTED_EXTENSIONS))}."
                )
            (raw_dir / filename).write_bytes(uploaded_file.read())


def build_grouped_reference_chunks_from_mentor(slug: str) -> dict[str, list[dict]]:
    grouped_chunks: dict[str, list[dict]] = {mode: [] for mode in MODE_DEFINITIONS}
    for mode, config in REFERENCE_UPLOAD_GROUPS.items():
        raw_dir = mode_dir_for_mentor(slug, mode) / "raw"
        if not raw_dir.exists():
            continue
        for path in sorted(raw_dir.iterdir(), key=lambda item: item.name.lower()):
            if not path.is_file():
                continue
            if not is_supported_filename(path.name):
                raise ValueError(
                    f"Unsupported file type for {config['label']}: {path.name}. "
                    f"Use one of: {', '.join(sorted(SUPPORTED_EXTENSIONS))}."
                )
            text = extract_text_from_upload(path.read_bytes(), path.name)
            chunks = chunk_text(text, mode)
            if chunks:
                grouped_chunks[mode].append({"source_file": path.name, "chunks": chunks})
    return grouped_chunks


def save_mentor_prompt_outputs(slug: str, artifacts: dict[str, PromptArtifact]) -> Path:
    directory = mentor_dir(slug)
    directory.mkdir(parents=True, exist_ok=True)
    for mode, artifact in artifacts.items():
        if artifact.record_count <= 0:
            continue
        mode_directory = mode_dir_for_mentor(slug, mode)
        mode_directory.mkdir(parents=True, exist_ok=True)
        (mode_directory / "prompt.txt").write_text(artifact.content, encoding="utf-8")
    (directory / "all_pi_style_prompts.txt").write_text(
        build_combined_prompt_text(artifacts),
        encoding="utf-8",
    )
    return directory


def render_home(**context: Any):
    prompt_mentors = list_prompt_mentors()
    selected_prompt_mentor = context.get("selected_prompt_mentor", "")
    if not selected_prompt_mentor:
        selected_prompt_mentor = request.args.get("prompt_mentor", "").strip().lower()
    if selected_prompt_mentor and not read_prompt_mentor(selected_prompt_mentor):
        selected_prompt_mentor = ""
    selected_prompt_mentor_profile = read_prompt_mentor(selected_prompt_mentor) if selected_prompt_mentor else None
    defaults = {
        "error": "",
        "feedback": "",
        "filename": "",
        "selected_type": "",
        "selected_mentor": DEFAULT_MENTOR_ID,
        "prompt_artifacts": [],
        "prompt_cards": [],
        "prompt_run_id": "",
        "prompt_run_location": "",
        "prompt_output_location": "",
        "prompt_download_urls": {},
        "prompt_message": "",
        "prompt_clean_url": url_for("home", prompt_mentor=selected_prompt_mentor) + "#prompt-library"
        if selected_prompt_mentor
        else url_for("home") + "#prompt-library",
        "prompt_mentors": prompt_mentors,
        "selected_prompt_mentor": selected_prompt_mentor,
        "selected_prompt_mentor_profile": selected_prompt_mentor_profile,
        "stored_prompt_files": stored_files_for_mentor(selected_prompt_mentor),
        "reset_on_refresh": False,
        "form": {
            "project_name": "",
            "source_type": "Research_Meeting_Minutes",
            "source_date": "",
            "mentor_mode": "",
        },
    }
    defaults.update(context)
    return render_template(
        "index.html",
        source_types=SOURCE_TYPE_TO_MENTOR_MODE,
        supported_extensions=", ".join(sorted(SUPPORTED_EXTENSIONS)),
        reference_upload_groups=REFERENCE_UPLOAD_GROUPS,
        mode_definitions=MODE_DEFINITIONS,
        upload_types=UPLOAD_TYPES,
        mentors=MENTORS,
        render_prompt_preview=render_prompt_preview,
        **defaults,
    )


@app.get("/")
def home():
    selected_type = request.args.get("type", "").strip().lower()
    if selected_type not in UPLOAD_TYPES:
        selected_type = ""
    selected_mentor = request.args.get("mentor", DEFAULT_MENTOR_ID).strip().lower()
    if selected_mentor not in MENTORS:
        selected_mentor = DEFAULT_MENTOR_ID
    selected_prompt_mentor = request.args.get("prompt_mentor", "").strip().lower()
    return render_home(
        selected_type=selected_type,
        selected_mentor=selected_mentor,
        selected_prompt_mentor=selected_prompt_mentor,
    )


@app.post("/generate-prompts")
def generate_prompts():
    try:
        prompt_mentor = resolve_prompt_mentor(
            selected_slug=request.form.get("selected_prompt_mentor", ""),
            new_name=request.form.get("prompt_mentor_name", ""),
        )
        save_uploaded_reference_files(prompt_mentor["slug"])
        grouped_chunks = build_grouped_reference_chunks_from_mentor(prompt_mentor["slug"])
    except ValueError as exc:
        return render_home(error=str(exc)), 400

    if not any(grouped_chunks[mode] for mode in grouped_chunks):
        return render_home(
            error="Please upload at least one reference material file.",
            selected_prompt_mentor=prompt_mentor["slug"],
        ), 400

    artifacts = build_mode_prompt_artifacts(grouped_chunks)
    llm_generated_prompts: dict[str, str] = {}
    for mode in MODE_DEFINITIONS:
        if artifacts[mode].record_count <= 0:
            continue
        mode_chunks = [
            str(chunk).strip()
            for group in grouped_chunks.get(mode, [])
            for chunk in group.get("chunks", [])
            if str(chunk).strip()
        ]
        llm_prompt = generate_llm_prompt_for_mode(
            mode,
            artifacts[mode].source_files,
            mode_chunks,
            extract_generated_prompt(artifacts[mode].content),
        )
        if llm_prompt:
            llm_generated_prompts[mode] = llm_prompt
    if llm_generated_prompts:
        artifacts = build_mode_prompt_artifacts(grouped_chunks, generated_prompts=llm_generated_prompts)
    project_name = prompt_mentor["name"]
    run_id = save_prompt_outputs(artifacts, project_name)
    mentor_output_location = save_mentor_prompt_outputs(prompt_mentor["slug"], artifacts)
    prompt_output_location = str(mentor_output_location)
    prompt_run_location = str(get_output_dir() / run_id)
    prompt_cards = [
        {
            "mode": mode,
            "label": artifacts[mode].label,
            "preview": compact_prompt_preview(artifacts[mode]),
        }
        for mode in MODE_DEFINITIONS
        if artifacts[mode].record_count > 0
    ]
    prompt_download_urls = {
        mode: f"/download/{run_id}/{mode}_prompt" for mode in MODE_DEFINITIONS
    }
    prompt_download_urls["all"] = f"/download/{run_id}/all_pi_style_prompts"

    response = Response(
        render_home(
            prompt_artifacts=[artifacts[mode] for mode in MODE_DEFINITIONS],
            prompt_cards=prompt_cards,
            prompt_run_id=run_id,
            prompt_run_location=prompt_run_location,
            prompt_output_location=prompt_output_location,
            prompt_download_urls=prompt_download_urls,
            prompt_message="PI Style Prompts Ready",
            selected_prompt_mentor=prompt_mentor["slug"],
            selected_prompt_mentor_profile=prompt_mentor,
            stored_prompt_files=stored_files_for_mentor(prompt_mentor["slug"]),
            prompt_clean_url=url_for("home", prompt_mentor=prompt_mentor["slug"]) + "#prompt-library",
            reset_on_refresh=True,
        )
    )
    response.headers["X-Prompt-Run-Id"] = run_id
    return response


def compact_prompt_preview(artifact: PromptArtifact) -> str:
    """Return one concise paragraph for the generated prompt result card."""
    lines = [line.strip() for line in artifact.content.splitlines() if line.strip()]
    generated_index = next(
        (
            index
            for index, line in enumerate(lines)
            if line == "Generated PI-style prompt:"
        ),
        -1,
    )
    if generated_index >= 0 and generated_index + 1 < len(lines):
        return truncate_preview(lines[generated_index + 1], limit=1100)

    instruction = next(
        (
            line
            for line in lines
            if line.startswith("You are reviewing")
        ),
        "Use the uploaded reference materials to generate direct, concrete PI-style feedback.",
    )
    reference_index = next(
        (
            index
            for index, line in enumerate(lines)
            if line == "Reference patterns extracted from uploaded raw materials:"
        ),
        -1,
    )
    reference = ""
    if reference_index >= 0:
        for line in lines[reference_index + 1 :]:
            if line.startswith("- ") and "<bullet>" not in line:
                reference = line[2:].strip()
                break
    prompt = instruction
    if reference:
        prompt = f"{instruction} Anchor the feedback around this uploaded pattern: {reference}"
    return truncate_preview(prompt, limit=420)


def split_prompt_preview_segments(preview: str) -> list[str]:
    """Split a long generated prompt preview into readable paragraph-like chunks."""
    normalized = " ".join(preview.split())
    if not normalized:
        return []
    sentences = [sentence.strip() for sentence in re.split(r"(?<=[.!?])\s+", normalized) if sentence.strip()]
    if len(sentences) <= 2:
        return [" ".join(sentences)]
    if len(sentences) <= 3:
        return sentences
    target_segments = min(MAX_PROMPT_PREVIEW_SEGMENTS, max(2, (len(sentences) + 2) // 3))
    base_size, remainder = divmod(len(sentences), target_segments)
    segments: list[str] = []
    cursor = 0
    for segment_index in range(target_segments):
        segment_size = base_size + (1 if segment_index < remainder else 0)
        segment = " ".join(sentences[cursor : cursor + segment_size])
        if segment:
            segments.append(segment)
        cursor += segment_size
    return segments


def render_prompt_preview(preview: str) -> Markup:
    """Render a prompt preview as safe, segmented text instead of keyword highlights."""
    segments = split_prompt_preview_segments(preview)
    if not segments:
        return Markup("")
    return Markup(
        "".join(
            f'<span class="prompt-segment">{html.escape(segment)}</span>'
            for segment in segments
        )
    )


def truncate_preview(text: str, limit: int) -> str:
    """Shorten preview text without cutting through a word or sentence when possible."""
    text = text.strip()
    if len(text) <= limit:
        return text
    sentence_cut = max(text.rfind(". ", 0, limit), text.rfind("? ", 0, limit), text.rfind("! ", 0, limit))
    if sentence_cut >= int(limit * 0.55):
        return text[: sentence_cut + 1].rstrip()
    word_cut = text.rfind(" ", 0, limit - 3)
    if word_cut <= 0:
        return text[: limit - 3].rstrip() + "..."
    return text[:word_cut].rstrip(" ,;:") + "..."


@app.post("/feedback")
def feedback():
    kind = request.form.get("content_type", "").strip().lower()
    mentor_id = request.form.get("mentor_id", DEFAULT_MENTOR_ID).strip().lower()
    focus = request.form.get("focus", "").strip()
    uploaded_file = request.files.get("file")

    if kind not in UPLOAD_TYPES:
        return render_home(error="Please choose what type of content you want reviewed."), 400
    if mentor_id not in MENTORS:
        return render_home(
            error="Please choose an available mentor.",
            selected_type=kind,
            selected_mentor=DEFAULT_MENTOR_ID,
        ), 400
    if not uploaded_file or not uploaded_file.filename:
        return render_home(
            error="Please choose a file to upload.",
            selected_type=kind,
            selected_mentor=mentor_id,
        ), 400

    filename = safe_uploaded_filename(uploaded_file)
    extension = Path(filename).suffix.lower()
    if extension not in UPLOAD_TYPES[kind]["extensions"]:
        allowed = ", ".join(sorted(UPLOAD_TYPES[kind]["extensions"]))
        return render_home(
            error=f"That file type is not supported for {UPLOAD_TYPES[kind]['label']}. Use: {allowed}.",
            selected_type=kind,
            selected_mentor=mentor_id,
        ), 400

    try:
        file_bytes = uploaded_file.read()
        extracted_text = extract_feedback_text(file_bytes, extension)
        prompt = build_feedback_prompt(kind, filename, extracted_text, focus, mentor_id)
        demo_feedback = (
            f"Demo feedback from {MENTORS[mentor_id]['name']} for {filename}\n\n"
            f"Your {UPLOAD_TYPES[kind]['label'].lower()} was uploaded and read successfully "
            f"({len(extracted_text)} characters extracted)."
        )
        model_feedback = call_model(prompt, demo_feedback=demo_feedback, mentor_id=mentor_id)
    except (requests.RequestException, ValueError) as exc:
        app.logger.exception("Feedback generation failed: %s", exc)
        return render_home(
            error="We couldn't generate feedback for that file.",
            selected_type=kind,
            selected_mentor=mentor_id,
        ), 502

    return render_home(
        feedback=model_feedback,
        filename=filename,
        selected_type=kind,
        selected_mentor=mentor_id,
        reset_on_refresh=True,
    )


def build_grouped_reference_chunks() -> dict[str, list[dict]]:
    grouped_chunks: dict[str, list[dict]] = {mode: [] for mode in MODE_DEFINITIONS}
    for mode, config in REFERENCE_UPLOAD_GROUPS.items():
        files = [file for file in request.files.getlist(config["field"]) if file and file.filename]
        for uploaded_file in files:
            filename = safe_uploaded_filename(uploaded_file)
            if not is_supported_filename(filename):
                raise ValueError(
                    f"Unsupported file type for {config['label']}: {filename}. "
                    f"Use one of: {', '.join(sorted(SUPPORTED_EXTENSIONS))}."
                )
            text = extract_text_from_upload(uploaded_file.read(), filename)
            chunks = chunk_text(text, mode)
            if not chunks:
                raise ValueError(f"No usable chunks were generated from {filename}.")
            grouped_chunks[mode].append({"source_file": filename, "chunks": chunks})
    return grouped_chunks


def safe_uploaded_filename(uploaded_file: FileStorage) -> str:
    original = uploaded_file.filename or "uploaded.txt"
    filename = secure_filename(original)
    return filename or Path(original).name


@app.post("/generate")
def generate():
    form = {
        "project_name": request.form.get("project_name", "").strip(),
        "source_type": request.form.get("source_type", "").strip(),
        "source_date": request.form.get("source_date", "").strip(),
        "mentor_mode": request.form.get("mentor_mode", "").strip(),
    }
    uploaded_file = request.files.get("file")

    if not form["project_name"]:
        return render_home(error="Project name is required.", form=form), 400

    if form["source_type"] not in SOURCE_TYPE_TO_MENTOR_MODE:
        return render_home(error="Choose one of the three source types.", form=form), 400

    if not form["mentor_mode"]:
        form["mentor_mode"] = default_mentor_mode(form["source_type"])

    if not uploaded_file or not uploaded_file.filename:
        return render_home(error="Please choose a raw material file.", form=form), 400

    filename = safe_uploaded_filename(uploaded_file)
    if not is_supported_filename(filename):
        return render_home(
            error=f"Unsupported file type. Use one of: {', '.join(sorted(SUPPORTED_EXTENSIONS))}.",
            form=form,
        ), 400

    try:
        file_bytes = uploaded_file.read()
        text = extract_text_from_upload(file_bytes, filename)
        chunks = chunk_text(text, form["source_type"])
        if not chunks:
            raise ValueError("No usable chunks were generated from this file.")
        records = build_records(
            chunks=chunks,
            project_name=form["project_name"],
            source_type=form["source_type"],
            source_date=form["source_date"],
            source_file=filename,
            mentor_mode=form["mentor_mode"],
        )
    except ValueError as exc:
        return render_home(error=str(exc), form=form), 400

    run_id = save_outputs(records, form["project_name"], form["source_type"])
    response = Response(
        render_home(
            records=records,
            record_count=len(records),
            run_id=run_id,
            download_urls={
                "jsonl": f"/download/{run_id}/jsonl",
                "json": f"/download/{run_id}/json",
                "md": f"/download/{run_id}/md",
            },
            form=form,
            reset_on_refresh=True,
        )
    )
    response.headers["X-Run-Id"] = run_id
    return response


def save_outputs(records: list[dict], project_name: str, source_type: str) -> str:
    run_id = f"{slugify(project_name, 4)}_{slugify(source_type, 4)}_{secrets.token_hex(4)}"
    run_dir = get_output_dir() / run_id
    run_dir.mkdir(parents=True, exist_ok=False)
    (run_dir / DOWNLOAD_FILENAMES["jsonl"]).write_text(records_to_jsonl(records), encoding="utf-8")
    (run_dir / DOWNLOAD_FILENAMES["json"]).write_text(records_to_pretty_json(records), encoding="utf-8")
    (run_dir / DOWNLOAD_FILENAMES["md"]).write_text(records_to_preview_markdown(records), encoding="utf-8")
    return run_id


def save_prompt_outputs(artifacts: dict[str, PromptArtifact], project_name: str) -> str:
    run_id = f"{slugify(project_name, 4)}_pi_prompts_{secrets.token_hex(4)}"
    run_dir = get_output_dir() / run_id
    run_dir.mkdir(parents=True, exist_ok=False)
    for artifact in artifacts.values():
        (run_dir / artifact.filename).write_text(artifact.content, encoding="utf-8")
    (run_dir / "all_pi_style_prompts.txt").write_text(
        build_combined_prompt_text(artifacts),
        encoding="utf-8",
    )
    return run_id


@app.get("/download/<run_id>/<kind>")
def download(run_id: str, kind: str):
    safe_run_id = secure_filename(run_id)
    if safe_run_id != run_id:
        abort(404)

    path = path_for_download(run_id, kind)
    if path is None or not path.exists():
        abort(404)

    mimetypes = {
        "jsonl": "application/jsonl; charset=utf-8",
        "json": "application/json; charset=utf-8",
        "md": "text/markdown; charset=utf-8",
    }
    mimetype = "text/plain; charset=utf-8" if path.suffix == ".txt" else mimetypes[kind]
    return send_file(
        path,
        mimetype=mimetype,
        as_attachment=True,
        download_name=path.name,
    )


def path_for_download(run_id: str, kind: str) -> Path | None:
    run_dir = get_output_dir() / run_id
    if kind in DOWNLOAD_FILENAMES:
        return run_dir / DOWNLOAD_FILENAMES[kind]
    if kind == "all_pi_style_prompts":
        return run_dir / "all_pi_style_prompts.txt"
    if kind.endswith("_prompt"):
        candidate = run_dir / f"{kind}.txt"
        if candidate.name in {f"{mode}_prompt.txt" for mode in MODE_DEFINITIONS}:
            return candidate
    return None


@app.errorhandler(413)
def file_too_large(_error: Exception):
    return render_home(error="The file is too large. Maximum upload size is 20 MB."), 413


@app.post("/api/generate")
def api_generate():
    body = request.get_json(silent=True) or {}
    prompt = body.get("prompt", "")
    prompt = prompt.strip() if isinstance(prompt, str) else ""
    if not prompt:
        return jsonify({"error": "Please enter a prompt."}), 400
    if len(prompt) > MAX_PROMPT_LENGTH:
        return jsonify({"error": "Prompt must be 4,000 characters or fewer."}), 400
    mentor_id = body.get("mentor_id", DEFAULT_MENTOR_ID)
    mentor_id = mentor_id.strip().lower() if isinstance(mentor_id, str) else ""
    if mentor_id not in MENTORS:
        return jsonify({"error": "Please choose an available mentor."}), 400
    try:
        return jsonify({"output": call_model(prompt, mentor_id=mentor_id)})
    except (requests.RequestException, ValueError) as exc:
        app.logger.exception("Model request failed: %s", exc)
        return jsonify({"error": "We couldn't reach the model."}), 502


if __name__ == "__main__":
    app.run(host="127.0.0.1", port=5000, debug=True)
